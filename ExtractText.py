# 11. Write Python chatbot to import: (20 points)
# First name initial A-F: Extract text from JPG images, see
# https://www.geeksforgeeks.org/how-to-extract-text-from-images-with-python/
# First name initial G-M: Extract text from PowerPoint files, search for "python extract text from pptx"
# First name initial N-Z: Extract text from Excel files, search for "python extract text from excel"

# First name initial: J => Extract text from PowerPoint files

import pyttsx3
from pptx import Presentation
import os
import time

# Initialize text-to-speech engine
engine = pyttsx3.init('sapi5')
voices = engine.getProperty('voices')
engine.setProperty('voice', voices[0].id)

# Text-to-Speech Function
def speak(text):
    engine.say(text)
    engine.runAndWait()

# Function to extract text from PowerPoint file
def extract_text_from_pptx(file_path):
    if not os.path.exists(file_path):
        print("The file does not exist.")
        speak("The file does not exist.")
        return
    
    presentation = Presentation(file_path)
    text_content = []

    # Loop through slides and extract text from each slide
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip() != "":
                text_content.append(shape.text.strip())

    # Join the extracted text
    return "\n".join(text_content)

# Greeting Function Based on Time of Day
def greetings():
    print("Hello, my name is Spartan. I can help you extract text from PowerPoint files.")
    speak("Hello, my name is Spartan. I can help you extract text from PowerPoint files.")
    print("Would you like me to extract text from a PowerPoint file?")
    speak("Would you like me to extract text from a PowerPoint file?")

# Main Command Function to interact with the user
def takeCommand():
    while True:
        print(" ")
        query = input("User: ")

        if 'yes' in query.lower() or 'again' in query.lower():
            print("Please provide the path to the PowerPoint file:")
            speak("Please provide the path to the PowerPoint file:")
            file_path = input("File Path: ")
            extracted_text = extract_text_from_pptx(file_path)

            if extracted_text:
                print("Here is the text extracted from your PowerPoint file:")
                print(extracted_text)
                speak("Here is the text extracted from your PowerPoint file.")
                speak(extracted_text)
            else:
                print("No text could be extracted from the PowerPoint file.")
                speak("No text could be extracted from the PowerPoint file.")
        
        elif 'exit' in query.lower() or 'bye' in query.lower():
            print("Goodbye, my friend. I will miss you.")
            speak("Goodbye, my friend. I will miss you.")
            break

# Startup and Greeting
time.sleep(2)
print('Initializing...')
time.sleep(2)
print('Spartan is preparing...')
time.sleep(2)
print('Environment is building...')
time.sleep(2)
greetings()
takeCommand()
